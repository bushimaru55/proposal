"""
PowerPoint生成機能
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import logging
from datetime import datetime

logger = logging.getLogger(__name__)


class PowerPointGenerator:
    """トークスクリプトからPowerPoint資料を生成"""
    
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
    
    def add_title_slide(self, company_name, generated_date=None):
        """タイトルスライド作成"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = f"{company_name}様\nご提案資料"
        
        if generated_date:
            subtitle.text = f"生成日: {generated_date}"
        else:
            subtitle.text = f"生成日: {datetime.now().strftime('%Y年%m月%d日')}"
        
        # タイトルのフォーマット
        title_frame = title.text_frame
        for paragraph in title_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(44)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0, 51, 102)
        
        return slide
    
    def add_section_slide(self, section_title, section_content):
        """セクションスライド作成"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = section_title
        
        # タイトルのフォーマット
        title_frame = title.text_frame
        for paragraph in title_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(32)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0, 51, 102)
        
        # コンテンツの追加
        text_frame = content.text_frame
        text_frame.clear()
        
        # セクションコンテンツを整形して追加
        paragraphs = section_content.split('\n')
        for para_text in paragraphs:
            if para_text.strip():
                p = text_frame.add_paragraph()
                p.text = para_text.strip()
                p.level = 0
                p.font.size = Pt(18)
        
        return slide
    
    def add_product_slide(self, product, proposal_angle):
        """商品紹介スライド作成"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = f"ご提案：{product.name}"
        
        # タイトルのフォーマット
        title_frame = title.text_frame
        for paragraph in title_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(32)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0, 51, 102)
        
        # コンテンツの追加
        text_frame = content.text_frame
        text_frame.clear()
        
        # 商品概要
        p = text_frame.add_paragraph()
        p.text = product.short_description
        p.font.size = Pt(16)
        p.space_after = Pt(12)
        
        # 提案角度
        if proposal_angle:
            p = text_frame.add_paragraph()
            p.text = f"【提案ポイント】"
            p.font.size = Pt(18)
            p.font.bold = True
            p.space_before = Pt(12)
            
            p = text_frame.add_paragraph()
            p.text = proposal_angle
            p.font.size = Pt(16)
        
        # 主要機能
        if product.key_features:
            p = text_frame.add_paragraph()
            p.text = "【主要機能】"
            p.font.size = Pt(18)
            p.font.bold = True
            p.space_before = Pt(12)
            
            for feature in product.key_features[:5]:  # 最大5つ
                p = text_frame.add_paragraph()
                feature_name = feature.get('name', '')
                feature_desc = feature.get('description', '')
                p.text = f"• {feature_name}: {feature_desc}"
                p.level = 1
                p.font.size = Pt(14)
        
        return slide
    
    def add_company_info_slide(self, company):
        """企業情報スライド作成"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = f"{company.company_name}様について"
        
        # コンテンツの追加
        text_frame = content.text_frame
        text_frame.clear()
        
        # 業界
        if company.industry:
            p = text_frame.add_paragraph()
            p.text = f"業界: {company.industry}"
            p.font.size = Pt(18)
        
        # 事業内容
        if company.business_description:
            p = text_frame.add_paragraph()
            p.text = "事業内容:"
            p.font.size = Pt(18)
            p.font.bold = True
            p.space_before = Pt(12)
            
            p = text_frame.add_paragraph()
            p.text = company.business_description
            p.font.size = Pt(16)
        
        # 推定される課題
        if company.pain_points:
            p = text_frame.add_paragraph()
            p.text = "推定される課題:"
            p.font.size = Pt(18)
            p.font.bold = True
            p.space_before = Pt(12)
            
            for pain_point in company.pain_points:
                p = text_frame.add_paragraph()
                p.text = f"• {pain_point}"
                p.level = 1
                p.font.size = Pt(16)
        
        return slide
    
    def generate_from_talk_script(self, talk_script):
        """トークスクリプトからPowerPoint生成"""
        
        # タイトルスライド
        self.add_title_slide(talk_script.company.company_name)
        
        # 企業情報スライド
        self.add_company_info_slide(talk_script.company)
        
        # セクション別スライド
        section_titles = {
            'opening': 'オープニング',
            'problem_identification': '課題の特定',
            'solution_proposal': 'ソリューション提案',
            'objection_handling': 'よくある懸念への対応',
            'closing': 'クロージング',
        }
        
        for section_name, content in talk_script.script_sections.items():
            if content and section_name in section_titles:
                self.add_section_slide(section_titles[section_name], content)
        
        # 提案商品スライド
        for link in talk_script.proposed_products.all().order_by('proposal_order'):
            proposal_angle = '\n'.join(link.matching_reasons) if link.matching_reasons else ''
            self.add_product_slide(link.product, proposal_angle)
        
        logger.info(f"PowerPoint生成完了: {len(self.prs.slides)}スライド")
        
        return self.prs
    
    def save(self, file_path):
        """PowerPointファイルを保存"""
        try:
            self.prs.save(file_path)
            logger.info(f"PowerPoint保存完了: {file_path}")
            return True
        except Exception as e:
            logger.error(f"PowerPoint保存エラー: {e}")
            return False

